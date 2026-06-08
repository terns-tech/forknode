#!/usr/bin/env python3
"""Generate EN/DE Canva-compatible PPTX pitch decks for HackPune 2026."""

import collections
import collections.abc

collections.Container = collections.abc.Container
collections.Mapping = collections.abc.Mapping
collections.MutableMapping = collections.abc.MutableMapping
collections.Sequence = collections.abc.Sequence

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

OUT = Path(__file__).parent
P = RGBColor(98, 27, 255)
L = RGBColor(231, 215, 255)
B = RGBColor(5, 5, 5)
G = RGBColor(242, 242, 242)
M = RGBColor(115, 115, 115)
W = RGBColor(255, 255, 255)

LANG = {
    "en": {
        "why": "Why We Are Doing This",
        "tag": "Where the world builders take flight.",
        "contact": "hello@terns.tech · terns.tech",
        "master": (
            "German companies need a low-risk way to access India's applied AI talent. "
            "Indian builders need real problems from real companies. HackPune connects "
            "8 German challenge owners with 160 curated Indian hackers in 24 hours."
        ),
        "stats": [
            ("109k", "unfilled IT roles in Germany"),
            ("85%", "German firms report IT shortage"),
            ("2.55M", "STEM graduates/year in India"),
            ("160", "Indian hackers in Pune"),
        ],
        "bridge": [
            "German startups",
            "HackPune Pune",
            "160 Indian builders",
            "Prototypes + hiring",
        ],
        "process": [
            "German challenge",
            "Curated teams",
            "24h build",
            "Pitch + recruit",
        ],
        "figs": [
            ("160", "Indian hackers"),
            ("8", "German challenges"),
            ("24h", "prototype sprint"),
            ("8", "winners TBD"),
        ],
        "sources": "Selected References",
        "cta": "Next Step",
        "stats_title": "Market Reality: The Gap Is Structural",
        "stats_note": (
            "HackPune is demand validation first — not logistics-first event planning."
        ),
        "process_title": "How HackPune Works",
        "process_note": (
            "All 8 challenges come from German startups and companies. "
            "Prizes are not fixed yet — to be decided with each partner."
        ),
        "value_title": "What Each Stakeholder Gets",
        "value_note": "Outcomes: prototype quality, talent access, credibility, community.",
        "timeline_title": "Timeline & Decision Points",
        "timeline_note": (
            "Critical validation: confirm German company interest before venue and media scale-up."
        ),
        "potential_title": "Potential Collaborations",
        "potential_note": "No partnerships confirmed yet — targets for outreach.",
    },
    "de": {
        "why": "Warum Wir Das Tun",
        "tag": "Where the world builders take flight.",
        "contact": "hello@terns.tech · terns.tech",
        "master": (
            "Deutsche Unternehmen brauchen einen risikoarmen Zugang zu angewandtem "
            "KI-Talent in Indien. Indische Builder brauchen echte Unternehmensprobleme. "
            "HackPune verbindet 8 deutsche Challenge-Owner mit 160 kuratierten "
            "indischen Hackern in 24 Stunden."
        ),
        "stats": [
            ("109k", "unbesetzte IT-Stellen in Deutschland"),
            ("85%", "Firmen sehen IT-Fachkraeftemangel"),
            ("2,55M", "STEM-Absolventen/Jahr in Indien"),
            ("160", "indische Hacker in Pune"),
        ],
        "bridge": [
            "Deutsche Startups",
            "HackPune Pune",
            "160 indische Builder",
            "Prototypen + Recruiting",
        ],
        "process": [
            "Deutsche Challenge",
            "Kuratierte Teams",
            "24h Build",
            "Pitch + Recruiting",
        ],
        "figs": [
            ("160", "indische Hacker"),
            ("8", "deutsche Challenges"),
            ("24h", "Prototyp-Sprint"),
            ("8", "Gewinner TBD"),
        ],
        "sources": "Ausgewaehlte Quellen",
        "cta": "Naechster Schritt",
        "stats_title": "Marktrealitaet: Die Luecke Ist Strukturell",
        "stats_note": (
            "HackPune validiert Nachfrage zuerst — nicht Logistik zuerst."
        ),
        "process_title": "So Funktioniert HackPune",
        "process_note": (
            "Alle 8 Challenges kommen von deutschen Startups und Unternehmen. "
            "Preise sind noch nicht festgelegt."
        ),
        "value_title": "Was Stakeholder Bekommen",
        "value_note": "Ergebnisse: Prototypen, Talentzugang, Glaubwuerdigkeit, Community.",
        "timeline_title": "Timeline & Entscheidungen",
        "timeline_note": (
            "Kritische Validierung: deutsches Unternehmensinteresse vor Venue- und Media-Skalierung."
        ),
        "potential_title": "Potenzielle Kooperationen",
        "potential_note": "Noch keine Partnerschaften bestaetigt — Outreach-Ziele.",
    },
}

BASE = {
    "company": {
        "en": [
            "HackPune 2026 — German Company Pitch",
            "24 hours of focused AI engineering on your real business problem — delivered by 160 Indian builders in Pune.",
            "For German startups and companies: CTOs, innovation leads, founders, Mittelstand.",
            "Book a discovery call and reserve one of 8 German challenge slots by end of July 2026.",
        ],
        "de": [
            "HackPune 2026 — Pitch Fuer Deutsche Unternehmen",
            "24 Stunden fokussiertes AI Engineering an Ihrem echten Problem — mit 160 indischen Buildern in Pune.",
            "Fuer deutsche Startups und Unternehmen: CTOs, Innovation, Gruender, Mittelstand.",
            "Discovery Call buchen und einen von 8 deutschen Challenge-Slots bis Ende Juli 2026 reservieren.",
        ],
    },
    "sponsor_partner": {
        "en": [
            "HackPune 2026 — Potential Partners & Collaborations",
            "Indo-German AI bridge — institutions, sponsors and ecosystem targets (not yet confirmed).",
            "For GINSEP, AHK, DWIH, DAAD, tool sponsors, venue and food partners.",
            "Express interest as a potential partner: co-brand, promote, fund or provide in-kind support.",
        ],
        "de": [
            "HackPune 2026 — Potenzielle Partner & Kooperationen",
            "Indo-German AI Bridge — Zielinstitutionen und Sponsoren (noch nicht bestaetigt).",
            "Fuer GINSEP, AHK, DWIH, DAAD, Tool-, Venue- und Food-Partner.",
            "Interesse als potenzieller Partner: Co-Branding, Promotion, Funding oder In-Kind.",
        ],
    },
    "participant_university": {
        "en": [
            "HackPune 2026 — Indian Participants & Universities",
            "160 Indian hackers. 8 German company challenges. 24 hours. Real prototypes.",
            "For Indian students, professionals and open-category builders; university promoters.",
            "Applications open soon — help us reach COEP, VIT Pune, PICT, SPPU and Pune tech communities.",
        ],
        "de": [
            "HackPune 2026 — Indische Teilnehmende & Universitaeten",
            "160 indische Hacker. 8 deutsche Company Challenges. 24 Stunden. Echte Prototypen.",
            "Fuer indische Studierende, Professionals und Open-Category Builder; Uni-Promotion.",
            "Bewerbungen bald offen — Reichweite ueber COEP, VIT Pune, PICT, SPPU und Pune Tech Community.",
        ],
    },
}

VALS = {
    "company": {
        "en": [
            ("Prototype ROI", "Multiple working prototypes on your exact German challenge."),
            ("Talent Signal", "Watch 160 Indian builders perform before hiring talks."),
            ("India Access", "Test India talent and market fit at ~€3.5k–5.5k per challenge slot."),
        ],
        "de": [
            ("Prototyp-ROI", "Mehrere Prototypen fuer Ihre deutsche Challenge."),
            ("Talent-Signal", "160 indische Builder vor Hiring-Gespraechen erleben."),
            ("India Access", "Indien-Talent testen ab ca. 3.500–5.500 EUR pro Challenge-Slot."),
        ],
    },
    "sponsor_partner": {
        "en": [
            ("Credibility", "Align with Indo-German innovation policy momentum."),
            ("Community", "Reach 160 Indian builders and Pune university networks."),
            ("Visibility", "Early-mover partner story across EN/DE media channels."),
        ],
        "de": [
            ("Glaubwuerdigkeit", "An Indo-German Innovation Momentum andocken."),
            ("Community", "160 indische Builder und Pune-Uni-Netzwerke erreichen."),
            ("Sichtbarkeit", "Early-Mover-Partner-Story in EN/DE Media."),
        ],
    },
    "participant_university": {
        "en": [
            ("Real Work", "8 challenges from German startups — no toy problems."),
            ("Career Access", "Meet German mentors, judges and recruiters in Pune."),
            ("Community", "Join Terns builder network beyond the 24h event."),
        ],
        "de": [
            ("Echte Arbeit", "8 Challenges von deutschen Startups — keine Toy Tasks."),
            ("Karriere", "Deutsche Mentoren, Judges und Recruiter in Pune treffen."),
            ("Community", "Terns Builder Network ueber das Event hinaus."),
        ],
    },
}

POTENTIAL = {
    "en": [
        "GINSEP — startup network (BMWK-backed)",
        "AHK / IGCC — Indo-German business chamber",
        "DWIH + DAAD — InnoXchange deep-tech programme",
        "Universities — COEP, SPPU, VIT Pune, PICT",
        "Tool sponsors — Cursor, OpenAI, AWS, Google Cloud",
    ],
    "de": [
        "GINSEP — Startup-Netzwerk (BMWK)",
        "AHK / IGCC — Deutsch-Indische Handelskammer",
        "DWIH + DAAD — InnoXchange Deep-Tech-Programm",
        "Universitaeten — COEP, SPPU, VIT Pune, PICT",
        "Tool-Sponsoren — Cursor, OpenAI, AWS, Google Cloud",
    ],
}

REFS = [
    "Bitkom 2025: 109k unfilled IT roles; 85% shortage; 79% expect worsening",
    "DWIH InnoXchange 2026: DST + DAAD bilateral deep-tech exchange",
    "PM India Jan 2026: bilateral trade > USD 50B; AI & innovation roadmap",
    "GINSEP: 1,500+ startups; 70+ market entries; BMWK-backed (potential partner)",
    "Inc42/Zinnov: Pune 360+ GCCs; top emerging hub; 90k graduates/year",
    "NITI/NASSCOM: India AI demand 1.25M+ by 2026–27; skills-based hiring rising",
    "ChefTreff 2026: proven 24h real-challenge format (reference model)",
    "terns.tech: HackPune 3–4 Oct 2026; 160 Indian hackers; 8 German challenges",
]


def tx(slide, x, y, w, h, text, size=24, bold=False, color=B, align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.margin_left = 0
    tf.margin_right = 0
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = "Arial"
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    if align:
        p.alignment = align
    return box


def logo(slide):
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.35), Inches(0.28), Inches(0.22), Inches(0.22))
    c.fill.solid()
    c.fill.fore_color.rgb = P
    c.line.color.rgb = P
    tx(slide, 0.65, 0.25, 1.5, 0.25, "TERNS", 13, True, B)


def bg(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    rect.fill.solid()
    rect.fill.fore_color.rgb = G
    rect.line.fill.background()
    logo(slide)
    return slide


def pill(slide, x, y, w, h, text, fill=L):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = fill
    tx(slide, x + 0.15, y + 0.09, w - 0.3, h - 0.1, text, 13, True, B, PP_ALIGN.CENTER)


def title_slide(prs, deck, lang):
    slide = bg(prs)
    t, sub, aud, _ = BASE[deck][lang]
    tx(slide, 0.55, 1.12, 5.6, 1.25, t, 34, True, B)
    tx(slide, 0.57, 2.55, 5.2, 0.9, sub, 18)
    tx(slide, 0.57, 3.55, 4.9, 0.5, aud, 13, False, M)
    panel = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.25), Inches(0), Inches(6.1), Inches(7.5))
    panel.fill.solid()
    panel.fill.fore_color.rgb = W
    panel.line.fill.background()
    for i, (num, lab) in enumerate(LANG[lang]["figs"]):
        x = 6.75 + (i % 2) * 2.35
        y = 1.0 + (i // 2) * 1.55
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(2.05), Inches(1.15))
        card.fill.solid()
        card.fill.fore_color.rgb = L if i == 0 else G
        card.line.color.rgb = W
        tx(slide, x + 0.15, y + 0.18, 1.25, 0.35, num, 28, True, P)
        tx(slide, x + 0.15, y + 0.68, 1.65, 0.3, lab, 10)
    tx(slide, 6.85, 5.0, 4.2, 0.55, LANG[lang]["tag"], 21, True)
    tx(slide, 0.55, 6.85, 3.5, 0.3, LANG[lang]["contact"], 9)


def why_slide(prs, lang):
    slide = bg(prs)
    tx(slide, 0.55, 0.8, 7, 0.5, LANG[lang]["why"], 30, True)
    tx(slide, 0.6, 1.45, 7.4, 1.1, LANG[lang]["master"], 22)
    xs = [0.75, 3.4, 6.1, 8.8]
    for i, lab in enumerate(LANG[lang]["bridge"]):
        pill(slide, xs[i], 3.3, 2.0, 0.85, lab, P if i == 1 else W)
        if i < 3:
            line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(xs[i] + 2.02), Inches(3.72), Inches(xs[i + 1] - 0.1), Inches(3.72))
            line.line.color.rgb = P
            line.line.width = Pt(2)
    tx(slide, 0.75, 5.15, 9.9, 0.55, "Problem → Bridge → Outcome", 26, True, P, PP_ALIGN.CENTER)
    tx(slide, 1.0, 5.9, 8.5, 0.55, LANG[lang]["stats_note"], 15, False, B, PP_ALIGN.CENTER)


def stats_slide(prs, lang):
    slide = bg(prs)
    tx(slide, 0.55, 0.78, 8, 0.5, LANG[lang]["stats_title"], 29, True)
    for i, (num, lab) in enumerate(LANG[lang]["stats"]):
        x = 0.75 + (i % 2) * 5.1
        y = 1.65 + (i // 2) * 1.75
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(4.45), Inches(1.25))
        card.fill.solid()
        card.fill.fore_color.rgb = W
        card.line.color.rgb = L
        tx(slide, x + 0.25, y + 0.22, 1.25, 0.42, num, 31, True, P)
        tx(slide, x + 1.55, y + 0.32, 2.65, 0.55, lab, 16)


def participant_mix_slide(prs, lang):
    slide = bg(prs)
    title = "160 Indian Hackers — Participant Mix" if lang == "en" else "160 Indische Hacker — Teilnehmer-Mix"
    tx(slide, 0.55, 0.78, 9, 0.5, title, 30, True)
    mix = [
        ("40%", "Students", "COEP, VIT Pune, PICT, SPPU, IIT Bombay"),
        ("40%", "Professional workers", "AI/ML engineers, developers, PMs"),
        ("20%", "Anyone", "Founders, designers, strategists, indie builders"),
    ]
    if lang == "de":
        mix = [
            ("40%", "Studierende", "COEP, VIT Pune, PICT, SPPU, IIT Bombay"),
            ("40%", "Berufstaetige", "AI/ML Engineers, Developer, PMs"),
            ("20%", "Alle", "Gruender, Designer, Strategen, Indie Builder"),
        ]
    for i, (pct, head, sub) in enumerate(mix):
        x = 0.85 + i * 3.65
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.9), Inches(3.1), Inches(2.4))
        card.fill.solid()
        card.fill.fore_color.rgb = W
        card.line.color.rgb = L
        tx(slide, x + 0.25, 2.15, 1.0, 0.45, pct, 34, True, P)
        tx(slide, x + 0.25, 2.75, 2.5, 0.35, head, 20, True)
        tx(slide, x + 0.25, 3.25, 2.55, 0.8, sub, 13, False, M)
    note = (
        "All participants are Indian residents. Curated applications; ~180 accepted, 160 seated."
        if lang == "en"
        else "Alle Teilnehmenden aus Indien. Kuratiert; ca. 180 angenommen, 160 Plaetze."
    )
    tx(slide, 1.0, 5.2, 9.0, 0.55, note, 16, True, B, PP_ALIGN.CENTER)


def process_slide(prs, lang):
    slide = bg(prs)
    tx(slide, 0.55, 0.78, 8, 0.5, LANG[lang]["process_title"], 30, True)
    for i, lab in enumerate(LANG[lang]["process"]):
        x = 0.85 + i * 2.75
        circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(2.1), Inches(1.15), Inches(1.15))
        circ.fill.solid()
        circ.fill.fore_color.rgb = P if i == 0 else L
        circ.line.color.rgb = P
        tx(slide, x + 0.37, 2.39, 0.4, 0.3, str(i + 1), 22, True, W if i == 0 else P, PP_ALIGN.CENTER)
        tx(slide, x - 0.35, 3.48, 1.85, 0.5, lab, 16, True, B, PP_ALIGN.CENTER)
        if i < 3:
            line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x + 1.25), Inches(2.67), Inches(x + 2.55), Inches(2.67))
            line.line.color.rgb = P
            line.line.width = Pt(2)
    tx(slide, 1.0, 5.0, 9.2, 0.8, LANG[lang]["process_note"], 17, False, B, PP_ALIGN.CENTER)


def value_slide(prs, deck, lang):
    slide = bg(prs)
    tx(slide, 0.55, 0.78, 8, 0.5, LANG[lang]["value_title"], 30, True)
    for i, (head, body) in enumerate(VALS[deck][lang]):
        x = 0.8 + i * 3.65
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.8), Inches(3.1), Inches(2.55))
        card.fill.solid()
        card.fill.fore_color.rgb = W
        card.line.color.rgb = L
        tx(slide, x + 0.25, 2.05, 2.3, 0.35, head, 20, True, P)
        tx(slide, x + 0.25, 2.65, 2.55, 1.0, body, 15)
    tx(slide, 1.15, 5.35, 8.8, 0.55, LANG[lang]["value_note"], 17, True, B, PP_ALIGN.CENTER)


def potential_slide(prs, lang):
    slide = bg(prs)
    tx(slide, 0.55, 0.78, 8, 0.5, LANG[lang]["potential_title"], 30, True)
    for i, item in enumerate(POTENTIAL[lang]):
        y = 1.55 + i * 0.72
        pill(slide, 0.95, y, 10.5, 0.55, item, L if i % 2 else W)
    tx(slide, 1.0, 5.35, 9.0, 0.55, LANG[lang]["potential_note"], 16, True, P, PP_ALIGN.CENTER)


def economics_slide(prs, lang):
    slide = bg(prs)
    title = "Economics (EUR)" if lang == "en" else "Oekonomie (EUR)"
    tx(slide, 0.55, 0.78, 8, 0.5, title, 30, True)
    rows = [
        ("Challenge slot", "€3,500 – €5,500", "8 German companies"),
        ("Title sponsor", "€5,500", "1 slot — branding + keynote"),
        ("Founder sponsor", "€100 – €500", "Community supporters"),
        ("Institutional", "€1,000 – €2,000", "Potential ecosystem partners"),
        ("Prizes", "TBD", "Decided with each challenge partner"),
        ("Operating budget", "~€12,000 – €15,000", "Venue, food, AV, ops"),
    ]
    if lang == "de":
        rows = [
            ("Challenge Slot", "3.500 – 5.500 EUR", "8 deutsche Unternehmen"),
            ("Title Sponsor", "5.500 EUR", "1 Slot — Branding + Keynote"),
            ("Founder Sponsor", "100 – 500 EUR", "Community Supporter"),
            ("Institutionell", "1.000 – 2.000 EUR", "Potenzielle Ecosystem Partner"),
            ("Preise", "TBD", "Mit jedem Challenge Partner"),
            ("Betriebsbudget", "ca. 12.000 – 15.000 EUR", "Venue, Food, AV, Ops"),
        ]
    for i, (a, b, c) in enumerate(rows):
        y = 1.45 + i * 0.72
        tx(slide, 0.95, y, 2.2, 0.35, a, 14, True, P)
        tx(slide, 3.2, y, 2.0, 0.35, b, 14, True, B)
        tx(slide, 5.4, y, 5.5, 0.35, c, 13, False, M)


def timeline_slide(prs, lang):
    slide = bg(prs)
    tx(slide, 0.55, 0.78, 7, 0.5, LANG[lang]["timeline_title"], 30, True)
    items = [
        ("Jun", "Partner outreach"),
        ("Jul", "8 German challenges"),
        ("Aug", "160 hackers selected"),
        ("Sep", "Dry run"),
        ("Oct 3-4", "HackPune"),
    ]
    if lang == "de":
        items = [
            ("Jun", "Partner Outreach"),
            ("Jul", "8 deutsche Challenges"),
            ("Aug", "160 Hacker selektiert"),
            ("Sep", "Dry Run"),
            ("Oct 3-4", "HackPune"),
        ]
    for i, (month, lab) in enumerate(items):
        x = 0.8 + i * 2.25
        tx(slide, x, 2.15, 1.2, 0.3, month, 18, True, P, PP_ALIGN.CENTER)
        pill(slide, x - 0.2, 2.75, 1.55, 0.55, lab, W)
        if i < 4:
            line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x + 1.35), Inches(3.02), Inches(x + 2.0), Inches(3.02))
            line.line.color.rgb = P
            line.line.width = Pt(2)
    tx(slide, 1.2, 5.15, 8.7, 0.65, LANG[lang]["timeline_note"], 17, True, B, PP_ALIGN.CENTER)


def cta_slide(prs, deck, lang):
    slide = bg(prs)
    tx(slide, 0.55, 0.85, 7, 0.5, LANG[lang]["cta"], 32, True)
    tx(slide, 0.75, 1.75, 7.8, 0.9, BASE[deck][lang][3], 24, True, P)
    steps = ["15-min call", "Fit check", "Confirm by Jul 2026"] if lang == "en" else ["15-Min Call", "Fit Check", "Bis Jul 2026"]
    for i, step in enumerate(steps):
        pill(slide, 0.95 + i * 3.15, 3.45, 2.55, 0.75, step, L if i == 1 else W)
    tx(slide, 0.78, 5.75, 4, 0.4, LANG[lang]["contact"], 17, True)


def refs_slide(prs, lang):
    slide = bg(prs)
    tx(slide, 0.55, 0.75, 7, 0.5, LANG[lang]["sources"], 29, True)
    for i, ref in enumerate(REFS):
        tx(slide, 0.8, 1.45 + i * 0.55, 10.6, 0.35, "• " + ref, 11)


def build(deck, lang):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    title_slide(prs, deck, lang)
    why_slide(prs, lang)
    stats_slide(prs, lang)
    participant_mix_slide(prs, lang)
    process_slide(prs, lang)
    value_slide(prs, deck, lang)
    potential_slide(prs, lang)
    economics_slide(prs, lang)
    timeline_slide(prs, lang)
    cta_slide(prs, deck, lang)
    refs_slide(prs, lang)
    path = OUT / f"Terns_HackPune_2026_{deck}_{lang}.pptx"
    prs.save(path)
    return path


def main():
    paths = []
    for deck in BASE:
        for lang in ("en", "de"):
            paths.append(build(deck, lang))
            print(paths[-1])
    print(f"Generated {len(paths)} PPTX files.")


if __name__ == "__main__":
    main()
